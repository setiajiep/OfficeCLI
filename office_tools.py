#!/usr/bin/env python3
"""
OfficeCLI - Complete Office & PDF Suite for Python
Handles PDF conversion, text extraction, document merging/splitting, PDF stamping,
watermarking, zip/unzip, and creation of DOCX, Excel XLSX, PDF reports, and PPTX presentations.
Integrates with Telegram for optional instant file delivery.
"""

import os
import sys
import io
import time
import subprocess
from telegram_utils import prompt_send_to_telegram, send_file_to_telegram

# ---------------------------------------------------------
# PDF & Office Utilities
# ---------------------------------------------------------

def convert_to_pdf(input_path, output_dir=None):
    """Convert DOCX, XLSX, PPTX, HTML, TXT or a folder of images to PDF"""
    if not os.path.exists(input_path):
        print(f"❌ Path tidak ditemukan: {input_path}")
        return None
    
    if os.path.isdir(input_path):
        out_pdf = os.path.join(output_dir, f"{os.path.basename(input_path.rstrip('/\\\\'))}.pdf") if output_dir else None
        return folder_to_pdf(input_path, output_path=out_pdf)
        
    if output_dir is None:
        output_dir = os.path.dirname(os.path.abspath(input_path))
        
    cmd = ["soffice", "--headless", "--convert-to", "pdf", input_path, "--outdir", output_dir]
    try:
        res = subprocess.run(cmd, capture_output=True, text=True, check=True)
        filename = os.path.splitext(os.path.basename(input_path))[0] + ".pdf"
        out_pdf = os.path.join(output_dir, filename)
        if os.path.exists(out_pdf):
            print(f"✅ Berhasil konversi ke PDF: {out_pdf}")
            return out_pdf
    except Exception as e:
        print(f"❌ Gagal konversi PDF dengan LibreOffice: {e}")
        ext = os.path.splitext(input_path)[1].lower()
        if ext in ['.png', '.jpg', '.jpeg', '.webp', '.bmp']:
            out_name = os.path.splitext(os.path.basename(input_path))[0] + ".pdf"
            return images_to_pdf([input_path], output_path=os.path.join(output_dir, out_name))
    return None

def extract_text(file_path):
    """Extract text content from PDF, DOCX, PPTX, XLSX, CSV, TXT, JSON, MD"""
    if not os.path.exists(file_path):
        return "❌ File tidak ditemukan."
    
    ext = os.path.splitext(file_path)[1].lower()
    
    if ext == ".pdf":
        try:
            import pdfplumber
            text = ""
            with pdfplumber.open(file_path) as pdf:
                for page in pdf.pages:
                    text += (page.extract_text() or "") + "\n"
            return text.strip() or "Halaman PDF kosong atau berupa gambar."
        except Exception:
            try:
                res = subprocess.run(["pdftotext", file_path, "-"], capture_output=True, text=True)
                return res.stdout.strip()
            except Exception as e:
                return f"Error membaca PDF: {e}"

    elif ext in [".docx", ".doc"]:
        try:
            import docx
            doc = docx.Document(file_path)
            full_text = [p.text for p in doc.paragraphs if p.text.strip()]
            for table in doc.tables:
                for row in table.rows:
                    full_text.append(" | ".join([cell.text.strip() for cell in row.cells]))
            return "\n".join(full_text)
        except Exception as e:
            return f"Error membaca DOCX: {e}"

    elif ext in [".pptx", ".ppt"]:
        try:
            import pptx
            prs = pptx.Presentation(file_path)
            slides_text = []
            for i, slide in enumerate(prs.slides, 1):
                slides_text.append(f"--- Slide {i} ---")
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            if paragraph.text.strip():
                                slides_text.append(paragraph.text.strip())
            return "\n".join(slides_text)
        except Exception as e:
            return f"Error membaca PowerPoint: {e}"

    elif ext in [".xlsx", ".xls", ".csv"]:
        try:
            import pandas as pd
            if ext == ".csv":
                df = pd.read_csv(file_path)
                return df.to_string()
            else:
                sheets_dict = pd.read_excel(file_path, sheet_name=None)
                out_str = []
                for sheet_name, df in sheets_dict.items():
                    out_str.append(f"--- Sheet: {sheet_name} ---")
                    out_str.append(df.to_string())
                return "\n\n".join(out_str)
        except Exception as e:
            return f"Error membaca Spreadsheet: {e}"

    elif ext in [".txt", ".md", ".json", ".py", ".html", ".css", ".js", ".sh", ".yaml", ".yml"]:
        try:
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                return f.read()
        except Exception as e:
            return f"Error membaca file teks: {e}"

    return "Unsupported file format for text extraction."

def get_doc_info(file_path):
    """Retrieve metadata and structure info of a document"""
    if not os.path.exists(file_path):
        return {}
    
    ext = os.path.splitext(file_path)[1].lower()
    info = {
        "filename": os.path.basename(file_path),
        "size_kb": round(os.path.getsize(file_path) / 1024, 2),
        "format": ext.lstrip('.').upper()
    }

    if ext == ".pdf":
        try:
            from pypdf import PdfReader
            reader = PdfReader(file_path)
            info["pages"] = len(reader.pages)
        except Exception:
            pass
    elif ext in [".docx", ".doc"]:
        try:
            import docx
            doc = docx.Document(file_path)
            info["paragraphs"] = len(doc.paragraphs)
            info["tables"] = len(doc.tables)
        except Exception:
            pass
    elif ext in [".pptx", ".ppt"]:
        try:
            import pptx
            prs = pptx.Presentation(file_path)
            info["slides"] = len(prs.slides)
        except Exception:
            pass
    elif ext in [".xlsx", ".xls"]:
        try:
            import pandas as pd
            xl = pd.ExcelFile(file_path)
            info["sheets"] = xl.sheet_names
        except Exception:
            pass
    elif ext in [".png", ".jpg", ".jpeg", ".webp", ".bmp"]:
        try:
            from PIL import Image
            with Image.open(file_path) as img:
                info["dimensions"] = f"{img.width}x{img.height}"
                info["mode"] = img.mode
        except Exception:
            pass
    return info

def merge_pdfs(pdf_list, output_path=None):
    """Merge multiple PDF files into one single PDF document"""
    from pypdf import PdfWriter
    if not pdf_list:
        print("❌ Tidak ada file PDF yang diberikan.")
        return None

    valid_files = [f for f in pdf_list if os.path.exists(f) and f.lower().endswith(".pdf")]
    if not valid_files:
        print("❌ Tidak ada file PDF valid yang ditemukan.")
        return None

    if output_path is None:
        first_dir = os.path.dirname(os.path.abspath(valid_files[0]))
        output_path = os.path.join(first_dir, f"merged_{int(time.time())}.pdf")

    writer = PdfWriter()
    merged_count = 0
    for pdf in valid_files:
        try:
            writer.append(pdf)
            merged_count += 1
        except Exception as e:
            print(f"⚠️ Error melembagakan PDF '{pdf}': {e}")

    if merged_count == 0:
        return None

    with open(output_path, "wb") as f:
        writer.write(f)

    print(f"✅ Berhasil menggabungkan {merged_count} PDF: {output_path}")
    return output_path

def split_pdf(pdf_path, pages="1", output_path=None):
    """Extract specific page numbers or page ranges from a PDF (e.g. '1-3, 5, 8')"""
    from pypdf import PdfReader, PdfWriter
    if not os.path.exists(pdf_path):
        print(f"❌ File PDF tidak ditemukan: {pdf_path}")
        return None

    reader = PdfReader(pdf_path)
    total_pages = len(reader.pages)
    selected_indices = set()

    for part in str(pages).split(","):
        part = part.strip()
        if not part:
            continue
        if "-" in part:
            sp = part.split("-")
            try:
                s, e = int(sp[0]), int(sp[1])
                for p in range(max(1, s), min(total_pages, e) + 1):
                    selected_indices.add(p - 1)
            except ValueError:
                pass
        else:
            try:
                p = int(part)
                if 1 <= p <= total_pages:
                    selected_indices.add(p - 1)
            except ValueError:
                pass

    if not selected_indices:
        print("❌ Halaman tidak valid.")
        return None

    writer = PdfWriter()
    for idx in sorted(selected_indices):
        writer.add_page(reader.pages[idx])

    if output_path is None:
        base, _ = os.path.splitext(pdf_path)
        output_path = f"{base}_pages_{pages.replace(' ', '')}.pdf"

    with open(output_path, "wb") as f:
        writer.write(f)

    print(f"✅ Berhasil mengekstrak halaman ({pages}): {output_path}")
    return output_path

def stamp_image_on_pdf(pdf_path, img_path, output_path=None, page_num="-1", x=350, y=80, width=150, height=60, auto_nobg=True, rotation=0):
    """Stamp signature or logo image onto PDF page(s) with transparent background support"""
    from pypdf import PdfReader, PdfWriter
    from reportlab.pdfgen import canvas
    from reportlab.lib.utils import ImageReader
    from PIL import Image

    if not os.path.exists(pdf_path) or not os.path.exists(img_path):
        print("❌ File PDF atau Gambar tidak ditemukan.")
        return None

    if output_path is None:
        base, _ = os.path.splitext(pdf_path)
        output_path = f"{base}_signed.pdf"

    stamp_img_path = img_path
    if auto_nobg:
        try:
            with Image.open(img_path).convert("RGBA") as raw_img:
                import numpy as np
                data = np.array(raw_img)
                r, g, b, a = data[:,:,0], data[:,:,1], data[:,:,2], data[:,:,3]
                white_mask = (r > 210) & (g > 210) & (b > 210)
                data[:, :, 3][white_mask] = 0
                clean_img = Image.fromarray(data)
                
                temp_nobg = f"/tmp/stamp_temp_{int(time.time())}.png"
                clean_img.save(temp_nobg)
                stamp_img_path = temp_nobg
        except Exception:
            stamp_img_path = img_path

    reader = PdfReader(pdf_path)
    total_pages = len(reader.pages)
    str_p = str(page_num).lower().strip()
    target_indices = []

    if str_p in ["-1", "last"]:
        target_indices = [total_pages - 1]
    elif str_p in ["all", "*"]:
        target_indices = list(range(total_pages))
    else:
        for part in str_p.split(","):
            part = part.strip()
            try:
                val = int(part)
                if val == -1:
                    target_indices.append(total_pages - 1)
                elif 1 <= val <= total_pages:
                    target_indices.append(val - 1)
            except ValueError:
                pass

    if not target_indices:
        target_indices = [total_pages - 1]

    writer = PdfWriter()
    img_reader = ImageReader(stamp_img_path)

    for i, page in enumerate(reader.pages):
        if i in target_indices:
            page_width = float(page.mediabox.width)
            page_height = float(page.mediabox.height)

            calc_x, calc_y = x, y
            if isinstance(calc_x, str):
                pos = calc_x.lower()
                if pos in ["bottom-right", "br"]:
                    calc_x = page_width - width - 50
                    calc_y = 50
                elif pos in ["bottom-left", "bl"]:
                    calc_x = 50
                    calc_y = 50
                elif pos in ["top-right", "tr"]:
                    calc_x = page_width - width - 50
                    calc_y = page_height - height - 50
                elif pos in ["top-left", "tl"]:
                    calc_x = 50
                    calc_y = page_height - height - 50
                elif pos in ["center", "c"]:
                    calc_x = (page_width - width) / 2
                    calc_y = (page_height - height) / 2
                else:
                    calc_x = page_width - width - 50
                    calc_y = 50

            packet = io.BytesIO()
            can = canvas.Canvas(packet, pagesize=(page_width, page_height))
            
            if rotation != 0:
                can.saveState()
                can.translate(float(calc_x) + width/2, float(calc_y) + height/2)
                can.rotate(float(rotation))
                can.drawImage(img_reader, -width/2, -height/2, width=float(width), height=float(height), mask='auto')
                can.restoreState()
            else:
                can.drawImage(img_reader, float(calc_x), float(calc_y), width=float(width), height=float(height), mask='auto')
            
            can.save()
            packet.seek(0)
            overlay_pdf = PdfReader(packet)
            page.merge_page(overlay_pdf.pages[0])

        writer.add_page(page)

    with open(output_path, "wb") as f:
        writer.write(f)

    if stamp_img_path != img_path and os.path.exists(stamp_img_path):
        try:
            os.remove(stamp_img_path)
        except Exception:
            pass

    print(f"✅ Stempel tanda tangan berhasil ditambahkan: {output_path}")
    return output_path

def add_watermark_to_pdf(pdf_path, text="CONFIDENTIAL", output_path=None):
    """Add a diagonal text watermark to all pages of a PDF document"""
    from pypdf import PdfReader, PdfWriter
    from reportlab.pdfgen import canvas
    from reportlab.lib.colors import Color

    if not os.path.exists(pdf_path):
        print(f"❌ File PDF tidak ditemukan: {pdf_path}")
        return None

    if output_path is None:
        base, _ = os.path.splitext(pdf_path)
        output_path = f"{base}_watermarked.pdf"

    reader = PdfReader(pdf_path)
    writer = PdfWriter()

    for page in reader.pages:
        w = float(page.mediabox.width)
        h = float(page.mediabox.height)

        packet = io.BytesIO()
        can = canvas.Canvas(packet, pagesize=(w, h))
        can.saveState()
        can.setFillColor(Color(0.5, 0.5, 0.5, alpha=0.25))
        can.setFont("Helvetica-Bold", 48)
        can.translate(w / 2, h / 2)
        can.rotate(45)
        can.drawCentredString(0, 0, text)
        can.restoreState()
        can.save()

        packet.seek(0)
        watermark_pdf = PdfReader(packet)
        page.merge_page(watermark_pdf.pages[0])
        writer.add_page(page)

    with open(output_path, "wb") as f:
        writer.write(f)

    print(f"✅ Watermark teks berhasil ditambahkan: {output_path}")
    return output_path

def protect_pdf(pdf_path, password="123", output_path=None):
    """Password protect / encrypt a PDF file"""
    from pypdf import PdfReader, PdfWriter
    if not os.path.exists(pdf_path):
        print(f"❌ File PDF tidak ditemukan: {pdf_path}")
        return None

    if output_path is None:
        base, _ = os.path.splitext(pdf_path)
        output_path = f"{base}_protected.pdf"

    reader = PdfReader(pdf_path)
    writer = PdfWriter()
    for page in reader.pages:
        writer.add_page(page)

    writer.encrypt(password)
    with open(output_path, "wb") as f:
        writer.write(f)

    print(f"✅ PDF berhasil dilindungi kata sandi: {output_path}")
    return output_path

def unprotect_pdf(pdf_path, password="123", output_path=None):
    """Remove password encryption from a PDF file"""
    from pypdf import PdfReader, PdfWriter
    if not os.path.exists(pdf_path):
        print(f"❌ File PDF tidak ditemukan: {pdf_path}")
        return None

    if output_path is None:
        base, _ = os.path.splitext(pdf_path)
        output_path = f"{base}_unlocked.pdf"

    reader = PdfReader(pdf_path)
    if reader.is_encrypted:
        reader.decrypt(password)

    writer = PdfWriter()
    for page in reader.pages:
        writer.add_page(page)

    with open(output_path, "wb") as f:
        writer.write(f)

    print(f"✅ Kata sandi PDF berhasil dibuka: {output_path}")
    return output_path

def compress_pdf(pdf_path, output_path=None):
    """Compress PDF streams & reduce file size"""
    from pypdf import PdfReader, PdfWriter
    if not os.path.exists(pdf_path):
        print(f"❌ File PDF tidak ditemukan: {pdf_path}")
        return None

    if output_path is None:
        base, _ = os.path.splitext(pdf_path)
        output_path = f"{base}_compressed.pdf"

    reader = PdfReader(pdf_path)
    writer = PdfWriter()
    for page in reader.pages:
        page.compress_content_streams()
        writer.add_page(page)

    with open(output_path, "wb") as f:
        writer.write(f)

    orig_sz = os.path.getsize(pdf_path) / 1024
    new_sz = os.path.getsize(output_path) / 1024
    print(f"✅ PDF berhasil dikompres: {output_path} ({orig_sz:.1f} KB -> {new_sz:.1f} KB)")
    return output_path

def pdf_to_images(pdf_path, output_dir=None):
    """Convert each page of a PDF file into PNG images"""
    from PIL import Image
    if not os.path.exists(pdf_path):
        print(f"❌ File PDF tidak ditemukan: {pdf_path}")
        return []

    if output_dir is None:
        base = os.path.basename(pdf_path).replace(".pdf", "")
        output_dir = os.path.join(os.path.dirname(pdf_path), f"{base}_pages_img")
    os.makedirs(output_dir, exist_ok=True)

    out_images = []
    try:
        import pdfplumber
        with pdfplumber.open(pdf_path) as pdf:
            for i, page in enumerate(pdf.pages, 1):
                img = page.to_image(resolution=150)
                img_file = os.path.join(output_dir, f"page_{i}.png")
                img.save(img_file)
                out_images.append(img_file)
        print(f"✅ Berhasil mengkonversi PDF ke {len(out_images)} gambar di: {output_dir}")
        return out_images
    except Exception as e:
        print(f"❌ Gagal konversi PDF ke Gambar: {e}")
        return []

def images_to_pdf(img_list, output_path=None):
    """Combine multiple images or folders of images into a single PDF file"""
    from PIL import Image
    if isinstance(img_list, str):
        img_list = [img_list]
    
    valid_exts = ('.png', '.jpg', '.jpeg', '.webp', '.bmp')
    extracted_imgs = []
    
    for item in img_list:
        if not os.path.exists(item):
            continue
        if os.path.isdir(item):
            for root_dir, _, files in os.walk(item):
                for fname in sorted(files):
                    if fname.lower().endswith(valid_exts):
                        extracted_imgs.append(os.path.join(root_dir, fname))
        elif item.lower().endswith(valid_exts):
            extracted_imgs.append(item)
            
    extracted_imgs.sort()
    
    if not extracted_imgs:
        print("❌ Tidak ada file gambar yang valid.")
        return None

    if output_path is None:
        first_item = img_list[0]
        if os.path.isdir(first_item):
            folder_name = os.path.basename(first_item.rstrip('/\\'))
            output_path = os.path.join(first_item, f"{folder_name}.pdf")
        else:
            first_dir = os.path.dirname(os.path.abspath(first_item))
            output_path = os.path.join(first_dir, f"album_{int(time.time())}.pdf")

    pil_images = []
    for f in extracted_imgs:
        try:
            im = Image.open(f).convert("RGB")
            pil_images.append(im)
        except Exception:
            pass

    if pil_images:
        pil_images[0].save(output_path, save_all=True, append_images=pil_images[1:])
        print(f"✅ Berhasil menggabungkan {len(pil_images)} gambar ke PDF: {output_path}")
        return output_path
    return None

def folder_to_pdf(folder_path, output_path=None):
    """Combine all image files inside a folder into a single PDF document inside the folder or output path"""
    if not os.path.exists(folder_path):
        print(f"❌ Folder tidak ditemukan: {folder_path}")
        return None
    if not os.path.isdir(folder_path):
        print(f"❌ Path bukan folder: {folder_path}")
        return None
    return images_to_pdf([folder_path], output_path=output_path)


# ---------------------------------------------------------
# Document Generation (DOCX, XLSX, PDF, PPTX)
# ---------------------------------------------------------

def create_docx(output_path=None, title="Dokumen Resmi", subtitle="OfficeCLI Document Generator", paragraphs=None, table_data=None):
    """Create a professionally styled DOCX Word Document"""
    import docx
    from docx.shared import Inches, Pt, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.enum.table import WD_TABLE_ALIGNMENT

    if output_path is None:
        output_path = f"/root/Dokumen_{int(time.time())}.docx"

    doc = docx.Document()
    
    # Title
    p_title = doc.add_paragraph()
    p_title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run_title = p_title.add_run(title)
    run_title.font.name = 'Calibri'
    run_title.font.size = Pt(24)
    run_title.font.bold = True
    run_title.font.color.rgb = RGBColor(30, 60, 114)

    # Subtitle
    if subtitle:
        p_sub = doc.add_paragraph()
        p_sub.alignment = WD_ALIGN_PARAGRAPH.CENTER
        run_sub = p_sub.add_run(subtitle)
        run_sub.font.name = 'Calibri'
        run_sub.font.size = Pt(12)
        run_sub.font.italic = True
        run_sub.font.color.rgb = RGBColor(100, 100, 100)

    doc.add_paragraph() # Spacer

    # Paragraphs
    if paragraphs:
        if isinstance(paragraphs, str):
            paragraphs = [paragraphs]
        for text in paragraphs:
            p = doc.add_paragraph()
            p.paragraph_format.space_after = Pt(8)
            run = p.add_run(text)
            run.font.name = 'Calibri'
            run.font.size = Pt(11)
    else:
        p = doc.add_paragraph()
        run = p.add_run("Dokumen ini dibuat secara otomatis menggunakan OfficeCLI Suite.")
        run.font.name = 'Calibri'
        run.font.size = Pt(11)

    # Table
    if table_data and isinstance(table_data, list) and len(table_data) > 0:
        doc.add_paragraph()
        rows_cnt = len(table_data)
        cols_cnt = max(len(r) for r in table_data) if rows_cnt > 0 else 0

        table = doc.add_table(rows=rows_cnt, cols=cols_cnt)
        table.alignment = WD_TABLE_ALIGNMENT.CENTER
        table.style = 'Table Grid'

        for r_idx, row in enumerate(table_data):
            for c_idx, val in enumerate(row):
                cell = table.cell(r_idx, c_idx)
                cell.text = str(val)
                if r_idx == 0:
                    for p in cell.paragraphs:
                        for run in p.runs:
                            run.font.bold = True
                            run.font.color.rgb = RGBColor(255, 255, 255)

    doc.save(output_path)
    print(f"✅ Dokumen Word (DOCX) berhasil dibuat: {output_path}")
    return output_path

def create_excel(output_path=None, sheet_name="Data", headers=None, rows=None):
    """Create a styled Excel spreadsheet (.xlsx)"""
    import openpyxl
    from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
    from openpyxl.utils import get_column_letter

    if output_path is None:
        output_path = f"/root/Spreadsheet_{int(time.time())}.xlsx"

    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = sheet_name

    if headers is None:
        headers = ["No", "Item / Nama", "Kategori", "Jumlah", "Keterangan"]
    if rows is None:
        rows = [
            [1, "Project Alpha", "Development", 10, "Selesai"],
            [2, "Project Beta", "Design", 5, "In Progress"],
            [3, "Project Gamma", "Testing", 8, "Pending"]
        ]

    # Write Headers
    ws.append(headers)
    header_fill = PatternFill(start_color="1E3C72", end_color="1E3C72", fill_type="solid")
    header_font = Font(name="Calibri", size=11, bold=True, color="FFFFFF")
    align_center = Alignment(horizontal="center", vertical="center")

    for col in range(1, len(headers) + 1):
        cell = ws.cell(row=1, column=col)
        cell.fill = header_fill
        cell.font = header_font
        cell.alignment = align_center

    # Write Rows
    row_fill_alt = PatternFill(start_color="F2F5F9", end_color="F2F5F9", fill_type="solid")
    thin_border = Border(
        left=Side(style='thin', color='CCCCCC'),
        right=Side(style='thin', color='CCCCCC'),
        top=Side(style='thin', color='CCCCCC'),
        bottom=Side(style='thin', color='CCCCCC')
    )

    for r_idx, row_data in enumerate(rows, start=2):
        ws.append(row_data)
        use_alt = (r_idx % 2 == 1)
        for c_idx in range(1, len(row_data) + 1):
            cell = ws.cell(row=r_idx, column=c_idx)
            cell.border = thin_border
            if use_alt:
                cell.fill = row_fill_alt

    # Auto-adjust column widths
    for col in ws.columns:
        max_len = max(len(str(cell.value or '')) for cell in col)
        col_letter = get_column_letter(col[0].column)
        ws.column_dimensions[col_letter].width = max(max_len + 4, 12)

    wb.save(output_path)
    print(f"✅ Spreadsheet Excel (XLSX) berhasil dibuat: {output_path}")
    return output_path

def create_pdf(output_path=None, title="Laporan Resmi", content=None, headers=None):
    """Create a formatted PDF document using ReportLab"""
    from reportlab.lib.pagesizes import letter
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, HRFlowable
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib import colors

    if output_path is None:
        output_path = f"/root/Laporan_{int(time.time())}.pdf"

    doc = SimpleDocTemplate(output_path, pagesize=letter, rightMargin=40, leftMargin=40, topMargin=40, bottomMargin=40)
    styles = getSampleStyleSheet()

    title_style = ParagraphStyle(
        'DocTitle',
        parent=styles['Heading1'],
        fontName='Helvetica-Bold',
        fontSize=22,
        textColor=colors.HexColor('#1E3C72'),
        alignment=1, # Center
        spaceAfter=12
    )

    body_style = ParagraphStyle(
        'DocBody',
        parent=styles['BodyText'],
        fontName='Helvetica',
        fontSize=10,
        leading=14,
        spaceAfter=8
    )

    elements = []
    elements.append(Paragraph(title, title_style))
    elements.append(HRFlowable(width="100%", thickness=2, color=colors.HexColor('#1E3C72'), spaceAfter=15))

    if content:
        if isinstance(content, str):
            content = [content]
        for p in content:
            elements.append(Paragraph(p, body_style))
    else:
        elements.append(Paragraph("Laporan ini disusun secara otomatis oleh OfficeCLI PDF Engine.", body_style))

    elements.append(Spacer(1, 15))

    # Add default table if headers provided
    if headers and isinstance(headers, list):
        table_data = [headers]
        if isinstance(headers[0], list):
            table_data = headers

        t = Table(table_data, hAlign='CENTER')
        t.setStyle(TableStyle([
            ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#1E3C72')),
            ('TEXTCOLOR', (0, 0), (-1, 0), colors.white),
            ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
            ('BOTTOMPADDING', (0, 0), (-1, 0), 8),
            ('BACKGROUND', (0, 1), (-1, -1), colors.HexColor('#F8F9FA')),
            ('GRID', (0, 0), (-1, -1), 0.5, colors.HexColor('#CCCCCC')),
        ]))
        elements.append(t)

    doc.build(elements)
    print(f"✅ Dokumen PDF berhasil dibuat: {output_path}")
    return output_path

def create_pptx(output_path=None, title="Presentasi Utama", subtitle="Dibuat oleh OfficeCLI", slides_data=None):
    """Create a PowerPoint (.pptx) presentation"""
    import pptx
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor

    if output_path is None:
        output_path = f"/root/Presentasi_{int(time.time())}.pptx"

    prs = pptx.Presentation()

    # Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1]

    title_shape.text = title
    subtitle_shape.text = subtitle

    # Add Slides
    if slides_data and isinstance(slides_data, list):
        bullet_slide_layout = prs.slide_layouts[1]
        for s_info in slides_data:
            s = prs.slides.add_slide(bullet_slide_layout)
            shapes = s.shapes
            s_title = shapes.title
            s_body = shapes.placeholders[1]

            if isinstance(s_info, dict):
                s_title.text = s_info.get("title", "Slide Baru")
                tf = s_body.text_frame
                tf.word_wrap = True
                bullets = s_info.get("bullets", [])
                for i, b in enumerate(bullets):
                    if i == 0:
                        p = tf.paragraphs[0]
                        p.text = b
                    else:
                        p = tf.add_paragraph()
                        p.text = b
            else:
                s_title.text = str(s_info)

    prs.save(output_path)
    print(f"✅ Presentasi PowerPoint (PPTX) berhasil dibuat: {output_path}")
    return output_path

def create_zip(target_path, output_zip=None):
    """Compress file or directory into a ZIP archive"""
    import zipfile
    if not os.path.exists(target_path):
        print(f"❌ Target tidak ditemukan: {target_path}")
        return None
    if output_zip is None:
        base = os.path.basename(os.path.abspath(target_path))
        parent = os.path.dirname(os.path.abspath(target_path))
        output_zip = os.path.join(parent, f"{base}.zip")

    with zipfile.ZipFile(output_zip, 'w', zipfile.ZIP_DEFLATED) as zipf:
        if os.path.isdir(target_path):
            for root, dirs, files in os.walk(target_path):
                for file in files:
                    full_path = os.path.join(root, file)
                    rel_path = os.path.relpath(full_path, os.path.dirname(target_path))
                    zipf.write(full_path, rel_path)
        else:
            zipf.write(target_path, os.path.basename(target_path))

    print(f"✅ Berhasil membuat arsip ZIP: {output_zip}")
    return output_zip

def extract_zip(zip_path, extract_dir=None):
    """Extract a ZIP archive"""
    import zipfile
    if not os.path.exists(zip_path):
        print(f"❌ File ZIP tidak ditemukan: {zip_path}")
        return None
    if extract_dir is None:
        extract_dir = os.path.dirname(os.path.abspath(zip_path))

    with zipfile.ZipFile(zip_path, 'r') as zipf:
        zipf.extractall(extract_dir)

    print(f"✅ Berhasil mengekstrak ZIP ke: {extract_dir}")
    return extract_dir

def parse_float_arg(val, default):
    try:
        return float(val)
    except (ValueError, TypeError):
        return default

# ---------------------------------------------------------
# CLI Main Function
# ---------------------------------------------------------

if __name__ == "__main__":
    args = sys.argv[1:]
    send_tg = False

    # Check for --send-telegram or -t flag
    if "--send-telegram" in args:
        send_tg = True
        args.remove("--send-telegram")
    elif "-t" in args:
        send_tg = True
        args.remove("-t")

    if not args:
        print("""
🏢 OfficeCLI - Suite Manajement & Pembuat Dokumen
Usage: python3 office_tools.py <action> [args...] [--send-telegram / -t]

Tindakan yang tersedia:
  • convert_pdf <file>                  - Konversi Word/Excel/PPT ke PDF
  • extract_text <file>                 - Ekstrak teks dari PDF/Doc/Xls/Txt
  • info <file>                         - Tampilkan metadata & info dokumen
  • merge <out_pdf> <pdf1> <pdf2>...   - Gabungkan beberapa PDF
  • split <pdf> <halaman> [out_pdf]     - Ekstrak halaman (cth: "1-3,5")
  • stamp_pdf <pdf> <gambar> [out] [hal] [x] [y] [w] [h] - Cap TTD pada PDF
  • watermark_pdf <pdf> [teks] [out]   - Tambahkan watermark teks ke PDF
  • create_docx [out] [judul] [isi]     - Buat dokumen Word (.docx)
  • create_excel [out] [sheet] [headers] - Buat spreadsheet Excel (.xlsx)
  • create_pdf [out] [judul] [isi]      - Buat dokumen PDF (.pdf)
  • create_pptx [out] [judul] [sub]     - Buat presentasi PowerPoint (.pptx)
  • zip <target> [out_zip]              - Buat arsip ZIP
  • unzip <zip_file> [out_dir]          - Ekstrak arsip ZIP

Setiap pembuatan/perubahan dokumen akan menanyakan konfirmasi pengiriman ke Telegram.
Gunakan flag --send-telegram atau -t untuk pengiriman otomatis.
""")
        sys.exit(0)

    action = args[0]
    result_file = None

    if action == "convert_pdf" and len(args) > 1:
        result_file = convert_to_pdf(args[1])
    elif action == "extract_text" and len(args) > 1:
        print(extract_text(args[1]))
    elif action == "info" and len(args) > 1:
        import json
        print(json.dumps(get_doc_info(args[1]), indent=2))
    elif action == "merge" and len(args) > 2:
        out_file = args[1] if args[1] != "-" else None
        pdf_inputs = args[2:]
        result_file = merge_pdfs(pdf_inputs, output_path=out_file)
    elif action == "split" and len(args) > 2:
        pdf_file = args[1]
        pages_sel = args[2]
        out_file = args[3] if len(args) > 3 and args[3] != "-" else None
        result_file = split_pdf(pdf_file, pages=pages_sel, output_path=out_file)
    elif action == "stamp_pdf" and len(args) > 2:
        pdf_path = args[1]
        img_path = args[2]
        output_path = args[3] if len(args) > 3 and args[3] != "-" else None
        page_num = args[4] if len(args) > 4 else "-1"
        arg5 = args[5] if len(args) > 5 else "bottom-right"
        try:
            x = float(arg5)
        except ValueError:
            x = arg5
        y = parse_float_arg(args[6] if len(args) > 6 else None, 80)
        width = parse_float_arg(args[7] if len(args) > 7 else None, 150)
        height = parse_float_arg(args[8] if len(args) > 8 else None, 60)
        result_file = stamp_image_on_pdf(pdf_path, img_path, output_path, page_num, x, y, width, height)
    elif action == "watermark_pdf" and len(args) > 1:
        pdf_path = args[1]
        txt = args[2] if len(args) > 2 else "CONFIDENTIAL"
        out = args[3] if len(args) > 3 else None
        result_file = add_watermark_to_pdf(pdf_path, text=txt, output_path=out)
    elif action == "protect_pdf" and len(args) > 1:
        pdf_path = args[1]
        pwd = args[2] if len(args) > 2 else "123"
        out = args[3] if len(args) > 3 else None
        result_file = protect_pdf(pdf_path, password=pwd, output_path=out)
    elif action == "unprotect_pdf" and len(args) > 1:
        pdf_path = args[1]
        pwd = args[2] if len(args) > 2 else "123"
        out = args[3] if len(args) > 3 else None
        result_file = unprotect_pdf(pdf_path, password=pwd, output_path=out)
    elif action == "compress_pdf" and len(args) > 1:
        pdf_path = args[1]
        out = args[2] if len(args) > 2 else None
        result_file = compress_pdf(pdf_path, output_path=out)
    elif action == "pdf2img" and len(args) > 1:
        pdf_path = args[1]
        out_dir = args[2] if len(args) > 2 else None
        imgs = pdf_to_images(pdf_path, output_dir=out_dir)
        if imgs:
            result_file = imgs[0]
    elif action == "img2pdf" and len(args) > 1:
        img_inputs = args[1:]
        result_file = images_to_pdf(img_inputs)
    elif action == "folder_to_pdf" and len(args) > 1:
        folder_p = args[1]
        out_p = args[2] if len(args) > 2 else None
        result_file = folder_to_pdf(folder_p, output_path=out_p)

    elif action == "create_docx":
        out = args[1] if len(args) > 1 and args[1] != "-" else None
        title = args[2] if len(args) > 2 else "Dokumen Resmi"
        content = args[3] if len(args) > 3 else "Ini adalah isi dokumen Word yang dibuat secara otomatis."
        result_file = create_docx(output_path=out, title=title, paragraphs=[content])
    elif action == "create_excel":
        out = args[1] if len(args) > 1 and args[1] != "-" else None
        sheet = args[2] if len(args) > 2 else "Data"
        result_file = create_excel(output_path=out, sheet_name=sheet)
    elif action == "create_pdf":
        out = args[1] if len(args) > 1 and args[1] != "-" else None
        title = args[2] if len(args) > 2 else "Laporan Resmi"
        content = args[3] if len(args) > 3 else "Isi laporan PDF yang dibuat dengan OfficeCLI."
        result_file = create_pdf(output_path=out, title=title, content=[content])
    elif action == "create_pptx":
        out = args[1] if len(args) > 1 and args[1] != "-" else None
        title = args[2] if len(args) > 2 else "Presentasi OfficeCLI"
        sub = args[3] if len(args) > 3 else "Sub-judul Presentasi"
        result_file = create_pptx(output_path=out, title=title, subtitle=sub)
    elif action == "zip" and len(args) > 1:
        result_file = create_zip(args[1], args[2] if len(args) > 2 else None)
    elif action == "unzip" and len(args) > 1:
        extract_zip(args[1], args[2] if len(args) > 2 else None)

    # Prompt user or send to Telegram if a result file was generated
    if result_file and os.path.exists(result_file):
        prompt_send_to_telegram(result_file, force_send=send_tg)
